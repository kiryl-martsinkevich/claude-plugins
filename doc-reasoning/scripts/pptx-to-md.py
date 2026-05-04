#!/usr/bin/env python3
"""Extract text from PowerPoint (.pptx) files as markdown.

Usage: pptx-to-md.py <input.pptx> <output.md>
"""

import sys
from pathlib import Path


def extract_slides(filepath: str) -> list[dict]:
    try:
        from pptx import Presentation
    except ImportError:
        sys.stderr.write(
            "ERROR: python-pptx not installed. Run: pip install python-pptx\n"
        )
        sys.exit(1)

    prs = Presentation(filepath)
    slides = []

    for i, slide in enumerate(prs.slides, 1):
        texts = []
        notes = ""

        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        texts.append(text)

            if shape.has_table:
                table = shape.table
                rows = []
                for row in table.rows:
                    cells = [cell.text.strip() for cell in row.cells]
                    rows.append(" | ".join(cells))
                if rows:
                    texts.append("\n" + "\n".join(rows))

        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text.strip()

        if texts:
            slides.append({
                "num": i,
                "text": texts,
                "notes": notes,
            })

    return slides


def main():
    if len(sys.argv) != 3:
        print(f"Usage: {sys.argv[0]} <input.pptx> <output.md>", file=sys.stderr)
        sys.exit(1)

    input_path = sys.argv[1]
    output_path = sys.argv[2]

    if not Path(input_path).exists():
        print(f"ERROR: File not found: {input_path}", file=sys.stderr)
        sys.exit(1)

    slides = extract_slides(input_path)

    if not slides:
        print(f"WARNING: No text extracted from {input_path}", file=sys.stderr)
        Path(output_path).write_text(f"*(No text content extracted from {Path(input_path).name})*\n")
        return

    lines = [f"# {Path(input_path).stem}\n"]
    for slide in slides:
        lines.append(f"## Slide {slide['num']}\n")
        for text in slide["text"]:
            lines.append(text + "\n")
        if slide["notes"]:
            lines.append(f"> **Notes:** {slide['notes']}\n")

    Path(output_path).write_text("\n".join(lines))
    print(f"{output_path} ({len(slides)} slide(s))")


if __name__ == "__main__":
    main()
